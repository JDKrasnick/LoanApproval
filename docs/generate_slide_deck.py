from __future__ import annotations

from pathlib import Path
from typing import Iterable

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs"
ASSET_DIR = OUT_DIR / "deck_assets"
DECK_PATH = OUT_DIR / "LoanApproval_Product_Deck.pptx"


NAVY = RGBColor(15, 23, 42)
SLATE = RGBColor(71, 85, 105)
MUTED = RGBColor(100, 116, 139)
LINE = RGBColor(226, 232, 240)
SOFT = RGBColor(248, 250, 252)
BLUE = RGBColor(29, 78, 216)
CYAN = RGBColor(8, 145, 178)
EMERALD = RGBColor(5, 150, 105)
AMBER = RGBColor(217, 119, 6)
RED = RGBColor(220, 38, 38)
VIOLET = RGBColor(109, 40, 217)
WHITE = RGBColor(255, 255, 255)


def font(size: int = 24, bold: bool = False, color: RGBColor = NAVY):
    return {"size": Pt(size), "bold": bold, "color": color}


def set_run(run, size: int, bold: bool = False, color: RGBColor = NAVY):
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(slide, x, y, w, h, text: str = "", size: int = 18, bold: bool = False,
                color: RGBColor = NAVY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size, bold, color)
    return box


def add_slide_title(slide, title: str, kicker: str | None = None):
    if kicker:
        add_textbox(slide, Inches(0.6), Inches(0.42), Inches(4.2), Inches(0.28), kicker.upper(), 9, True, BLUE)
    add_textbox(slide, Inches(0.6), Inches(0.72), Inches(8.9), Inches(0.55), title, 26, True, NAVY)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(1.36), Inches(12.15), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def add_footer(slide, n: int):
    add_textbox(slide, Inches(0.6), Inches(7.03), Inches(2.2), Inches(0.22), "LoanApproval", 8, True, MUTED)
    add_textbox(slide, Inches(11.95), Inches(7.03), Inches(0.8), Inches(0.22), f"{n:02d}", 8, False, MUTED, PP_ALIGN.RIGHT)


def add_bullets(slide, x, y, w, h, items: Iterable[str], size: int = 15, color: RGBColor = SLATE):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(8)
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return box


def add_card(slide, x, y, w, h, title: str, body: str, accent: RGBColor = BLUE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Inches(0.05), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_textbox(slide, x + Inches(0.22), y + Inches(0.18), w - Inches(0.42), Inches(0.28), title, 13, True, NAVY)
    add_textbox(slide, x + Inches(0.22), y + Inches(0.52), w - Inches(0.42), h - Inches(0.62), body, 10, False, SLATE)


def add_image_placeholder(slide, x, y, w, h, label: str):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = SOFT
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1.1)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    set_run(run, 12, True, MUTED)
    return shape


def add_metric(slide, x, y, label: str, value: str, accent: RGBColor = BLUE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, Inches(2.45), Inches(0.86))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    add_textbox(slide, x + Inches(0.16), y + Inches(0.12), Inches(2.1), Inches(0.2), label.upper(), 7, True, MUTED)
    add_textbox(slide, x + Inches(0.16), y + Inches(0.38), Inches(2.1), Inches(0.3), value, 17, True, accent)


def load_font(size: int, bold: bool = False):
    candidates = [
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf" if bold else "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/Library/Fonts/Arial Bold.ttf" if bold else "/Library/Fonts/Arial.ttf",
    ]
    for path in candidates:
        if Path(path).exists():
            return ImageFont.truetype(path, size)
    return ImageFont.load_default()


def rounded_rect(draw, xy, radius, fill, outline=None, width=1):
    draw.rounded_rectangle(xy, radius=radius, fill=fill, outline=outline, width=width)


def make_architecture_image(path: Path):
    img = Image.new("RGB", (1600, 900), "#f8fafc")
    d = ImageDraw.Draw(img)
    title = load_font(46, True)
    body = load_font(25)
    small = load_font(20)
    d.text((70, 58), "LoanApproval System Architecture", font=title, fill="#0f172a")
    boxes = [
        (70, 185, 360, 325, "Borrower App", "Next.js application flow"),
        (455, 185, 745, 325, "FastAPI", "Validated REST endpoints"),
        (840, 185, 1130, 325, "Decisioning", "Rules + ML + LLM"),
        (1225, 185, 1515, 325, "Underwriter UI", "Portfolio + deal review"),
        (455, 500, 745, 640, "PostgreSQL", "Applications, decisions, logs"),
        (840, 500, 1130, 640, "Supabase Storage", "Supporting documents"),
    ]
    for x1, y1, x2, y2, h, b in boxes:
        rounded_rect(d, (x1, y1, x2, y2), 22, "#ffffff", "#cbd5e1", 3)
        d.rectangle((x1, y1, x1 + 12, y2), fill="#1d4ed8")
        d.text((x1 + 32, y1 + 32), h, font=body, fill="#0f172a")
        d.text((x1 + 32, y1 + 76), b, font=small, fill="#475569")
    arrows = [
        ((368, 255), (445, 255)),
        ((753, 255), (830, 255)),
        ((1138, 255), (1215, 255)),
        ((600, 333), (600, 490)),
        ((985, 333), (985, 490)),
    ]
    for start, end in arrows:
        d.line([start, end], fill="#64748b", width=5)
        ex, ey = end
        if start[1] == end[1]:
            d.polygon([(ex, ey), (ex - 16, ey - 10), (ex - 16, ey + 10)], fill="#64748b")
        else:
            d.polygon([(ex, ey), (ex - 10, ey - 16), (ex + 10, ey - 16)], fill="#64748b")
    d.text((70, 755), "Key pattern: structured application data and document events create a traceable credit file before decision automation runs.", font=small, fill="#334155")
    img.save(path)


def make_workflow_image(path: Path):
    img = Image.new("RGB", (1600, 900), "#ffffff")
    d = ImageDraw.Draw(img)
    title = load_font(44, True)
    body = load_font(24)
    small = load_font(19)
    d.text((70, 60), "End-to-End Loan Lifecycle", font=title, fill="#0f172a")
    stages = [
        ("Apply", "Business profile\nloan request\nfinancials"),
        ("Upload", "Statements\ntax returns\nbank statements"),
        ("Evaluate", "Core ratios\nindustry metrics\nhard stops"),
        ("Explain", "Plain-English\nrationale + tiers"),
        ("Manage", "Pipeline status\nfunding actions\naudit trail"),
    ]
    x = 90
    colors = ["#1d4ed8", "#0891b2", "#059669", "#d97706", "#6d28d9"]
    for i, (name, detail) in enumerate(stages):
        rounded_rect(d, (x, 230, x + 245, 540), 28, "#f8fafc", "#cbd5e1", 3)
        d.ellipse((x + 28, 260, x + 88, 320), fill=colors[i])
        d.text((x + 49, 274), str(i + 1), font=body, fill="#ffffff")
        d.text((x + 30, 350), name, font=body, fill="#0f172a")
        d.multiline_text((x + 30, 400), detail, font=small, fill="#475569", spacing=8)
        if i < len(stages) - 1:
            d.line((x + 252, 385, x + 310, 385), fill="#94a3b8", width=5)
            d.polygon([(x + 310, 385), (x + 292, 373), (x + 292, 397)], fill="#94a3b8")
        x += 305
    d.text((90, 685), "Outcome states: New → In Review → Approved → Funded, with Declined handled as a terminal branch.", font=small, fill="#334155")
    img.save(path)


def make_decision_chart(path: Path):
    img = Image.new("RGB", (1600, 900), "#ffffff")
    d = ImageDraw.Draw(img)
    title = load_font(44, True)
    label = load_font(24, True)
    small = load_font(20)
    d.text((70, 60), "Decisioning Model: Deterministic Policy Before Automation", font=title, fill="#0f172a")
    rounded_rect(d, (120, 190, 510, 350), 24, "#eff6ff", "#bfdbfe", 3)
    rounded_rect(d, (610, 190, 1000, 350), 24, "#ecfdf5", "#a7f3d0", 3)
    rounded_rect(d, (1100, 190, 1480, 350), 24, "#fff7ed", "#fed7aa", 3)
    d.text((155, 228), "Company Score", font=label, fill="#1d4ed8")
    d.text((155, 277), "5 weighted core metrics", font=small, fill="#475569")
    d.text((645, 228), "Industry Score", font=label, fill="#059669")
    d.text((645, 277), "5 track-specific metrics", font=small, fill="#475569")
    d.text((1135, 228), "Loan Bracket", font=label, fill="#d97706")
    d.text((1135, 277), "Thresholds by amount", font=small, fill="#475569")
    d.line((510, 270, 610, 270), fill="#94a3b8", width=5)
    d.line((1000, 270, 1100, 270), fill="#94a3b8", width=5)
    d.text((322, 456), "Final Score = 0.6 × Company + 0.4 × Industry", font=label, fill="#0f172a")
    d.rectangle((250, 550, 1350, 632), fill="#f1f5f9", outline="#cbd5e1", width=2)
    d.rectangle((250, 550, 685, 632), fill="#fee2e2")
    d.rectangle((685, 550, 1015, 632), fill="#fef3c7")
    d.rectangle((1015, 550, 1350, 632), fill="#dcfce7")
    d.text((408, 574), "Decline", font=small, fill="#991b1b")
    d.text((795, 574), "Manual Review", font=small, fill="#92400e")
    d.text((1132, 574), "Approve", font=small, fill="#166534")
    d.text((250, 666), "< HR lower threshold", font=small, fill="#475569")
    d.text((690, 666), "HR band", font=small, fill="#475569")
    d.text((1018, 666), "≥ approve threshold and no caps", font=small, fill="#475569")
    d.text((120, 770), "Hard stops override scoring: reject-tier DSCR, LTV, bankruptcies, or industry metrics can force decline or cap at review.", font=small, fill="#334155")
    img.save(path)


def make_assets():
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    make_architecture_image(ASSET_DIR / "architecture.png")
    make_workflow_image(ASSET_DIR / "workflow.png")
    make_decision_chart(ASSET_DIR / "decision_model.png")


def screenshot_path() -> Path | None:
    matches = sorted(ROOT.glob("Screenshot 2026-04-22*PM.png"))
    return matches[0] if matches else None


def add_picture_fit(slide, image_path: Path, x, y, w, h):
    with Image.open(image_path) as img:
        iw, ih = img.size
    scale = min(w / iw, h / ih)
    pic_w = int(iw * scale)
    pic_h = int(ih * scale)
    return slide.shapes.add_picture(str(image_path), x + (w - pic_w) / 2, y + (h - pic_h) / 2, width=pic_w, height=pic_h)


def build_deck():
    make_assets()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slides = []

    # 1
    slide = prs.slides.add_slide(blank)
    slides.append(slide)
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    add_textbox(slide, Inches(0.75), Inches(0.7), Inches(3.2), Inches(0.35), "CORNELL FINTECH CLUB · SPRING 2026", 10, True, RGBColor(147, 197, 253))
    add_textbox(slide, Inches(0.75), Inches(1.55), Inches(7.4), Inches(1.35), "LoanApproval", 46, True, WHITE)
    add_textbox(slide, Inches(0.78), Inches(2.78), Inches(7.2), Inches(0.8), "Full-stack lending platform for borrower intake, automated credit policy, document workflows, and underwriter review.", 18, False, RGBColor(203, 213, 225))
    add_metric(slide, Inches(0.82), Inches(4.38), "Lifecycle", "Apply → Fund", CYAN)
    add_metric(slide, Inches(3.55), Inches(4.38), "Decision modes", "Rules · ML · LLM", EMERALD)
    add_metric(slide, Inches(6.28), Inches(4.38), "Stack", "Next.js + FastAPI", BLUE)
    add_image_placeholder(slide, Inches(9.05), Inches(1.25), Inches(3.35), Inches(4.65), "{{ INSERT_HERO_PRODUCT_SCREENSHOT }}")
    add_footer(slide, 1)

    # 2
    slide = prs.slides.add_slide(blank)
    slides.append(slide)
    add_slide_title(slide, "Product At A Glance", "Core thesis")
    add_textbox(slide, Inches(0.72), Inches(1.76), Inches(5.25), Inches(0.85), "LoanApproval turns a fragmented SMB loan process into a structured credit file with a transparent policy decision.", 22, True, NAVY)
    add_bullets(slide, Inches(0.75), Inches(2.95), Inches(5.15), Inches(2.2), [
        "Borrowers complete a guided multi-step application.",
        "Underwriters review every deal in a portfolio dashboard.",
        "Rules engine produces deterministic approve, decline, or manual-review outcomes.",
        "ML scoring and LLM extraction are wired as risk and productivity layers.",
    ])
    add_card(slide, Inches(6.4), Inches(1.78), Inches(2.8), Inches(1.18), "Borrower", "Self-serve application, document upload, confirmation and status views.", BLUE)
    add_card(slide, Inches(9.45), Inches(1.78), Inches(2.8), Inches(1.18), "Underwriter", "Portfolio KPIs, pipeline, sortable application queue, detailed credit review.", EMERALD)
    add_card(slide, Inches(6.4), Inches(3.25), Inches(2.8), Inches(1.18), "Policy", "Weighted core ratios, industry tracks, amount-based thresholds, hard stops.", AMBER)
    add_card(slide, Inches(9.45), Inches(3.25), Inches(2.8), Inches(1.18), "Platform", "FastAPI, PostgreSQL, Supabase Storage, swappable OpenAI or Anthropic client.", CYAN)
    add_image_placeholder(slide, Inches(6.4), Inches(4.78), Inches(5.85), Inches(1.25), "{{ INSERT_PRODUCT_OVERVIEW_COMPOSITE }}")
    add_footer(slide, 2)

    # 3
    slide = prs.slides.add_slide(blank)
    slides.append(slide)
    add_slide_title(slide, "Loan Lifecycle", "Workflow")
    add_picture_fit(slide, ASSET_DIR / "workflow.png", Inches(0.68), Inches(1.58), Inches(11.95), Inches(4.75))
    add_footer(slide, 3)

    # 4
    slide = prs.slides.add_slide(blank)
    slides.append(slide)
    add_slide_title(slide, "Borrower Entry Point", "Frontend")
    add_textbox(slide, Inches(0.75), Inches(1.76), Inches(4.9), Inches(0.55), "The landing page positions the product around speed, transparency, and no hard credit pull.", 19, True, NAVY)
    add_bullets(slide, Inches(0.78), Inches(2.55), Inches(4.75), Inches(2.45), [
        "Primary CTA routes directly to `/apply/step/1`.",
        "Financing options include term loans, lines of credit, equipment finance, and bridge loans.",
        "Eligibility copy sets basic borrower expectations before intake begins.",
    ], 14)
    add_image_placeholder(slide, Inches(6.15), Inches(1.68), Inches(5.95), Inches(4.72), "{{ INSERT_SCREENSHOT: Landing page / financing options }}")
    add_footer(slide, 4)

    # 5
    slide = prs.slides.add_slide(blank)
    slides.append(slide)
    add_slide_title(slide, "Guided Application Intake", "Borrower flow")
    add_card(slide, Inches(0.72), Inches(1.78), Inches(3.68), Inches(1.1), "1. Business Profile", "Legal name, DBA, industry, structure, EIN, address, and years in operation.", BLUE)
    add_card(slide, Inches(0.72), Inches(3.08), Inches(3.68), Inches(1.1), "2. Loan Request", "Amount, purpose, optional use-of-funds detail, and desired term.", CYAN)
    add_card(slide, Inches(0.72), Inches(4.38), Inches(3.68), Inches(1.1), "3. Review & Submit", "Certified application payload creates borrower and application records.", EMERALD)
    add_image_placeholder(slide, Inches(5.05), Inches(1.68), Inches(6.95), Inches(4.82), "{{ INSERT_SCREENSHOT: Multi-step application form }}")
    add_footer(slide, 5)

    # 6
    slide = prs.slides.add_slide(blank)
    slides.append(slide)
    add_slide_title(slide, "Financial Data Capture", "Credit file")
    add_textbox(slide, Inches(0.75), Inches(1.72), Inches(4.95), Inches(0.58), "The borrower workflow collects financial statements as structured values before underwriting.", 19, True, NAVY)
    add_bullets(slide, Inches(0.78), Inches(2.54), Inches(4.9), Inches(2.28), [
        "Core fields: revenue, EBITDA, debt, assets, collateral, annual debt service, bankruptcies.",
        "Optional fields improve ratio coverage: current assets, liabilities, EBIT, interest expense.",
        "The UI calculates DSCR, LTV, and leverage while the user enters data.",
    ], 14)
    add_card(slide, Inches(6.12), Inches(1.72), Inches(2.65), Inches(1.0), "DSCR", "EBITDA / annual debt service", EMERALD)
    add_card(slide, Inches(9.0), Inches(1.72), Inches(2.65), Inches(1.0), "LTV", "loan amount / collateral", AMBER)
    add_card(slide, Inches(6.12), Inches(3.02), Inches(2.65), Inches(1.0), "Leverage", "existing debt / EBITDA", BLUE)
    add_card(slide, Inches(9.0), Inches(3.02), Inches(2.65), Inches(1.0), "Liquidity", "current assets / liabilities", CYAN)
    add_image_placeholder(slide, Inches(6.12), Inches(4.45), Inches(5.53), Inches(1.3), "{{ INSERT_SCREENSHOT: Financials step with calculated ratios }}")
    add_footer(slide, 6)

    # 7
    slide = prs.slides.add_slide(blank)
    slides.append(slide)
    add_slide_title(slide, "Document Workflow", "Evidence")
    add_textbox(slide, Inches(0.75), Inches(1.76), Inches(4.75), Inches(0.6), "Document collection supports the credit file without forcing everything into manual email threads.", 19, True, NAVY)
    add_bullets(slide, Inches(0.78), Inches(2.58), Inches(4.65), Inches(2.2), [
        "Required slots: financial statements, tax returns, and bank statements.",
        "Optional slots: AR aging and collateral documentation.",
        "Upload endpoint records metadata and storage path, then writes an audit event.",
    ], 14)
    add_image_placeholder(slide, Inches(6.0), Inches(1.68), Inches(5.98), Inches(4.82), "{{ INSERT_SCREENSHOT: Document upload drop zones }}")
    add_footer(slide, 7)

    # 8
    slide = prs.slides.add_slide(blank)
    slides.append(slide)
    add_slide_title(slide, "Backend Architecture", "Platform")
    add_picture_fit(slide, ASSET_DIR / "architecture.png", Inches(0.68), Inches(1.54), Inches(11.95), Inches(4.86))
    add_footer(slide, 8)

    # 9
    slide = prs.slides.add_slide(blank)
    slides.append(slide)
    add_slide_title(slide, "Credit Policy Engine", "Decisioning")
    add_picture_fit(slide, ASSET_DIR / "decision_model.png", Inches(0.68), Inches(1.48), Inches(7.5), Inches(4.85))
    add_bullets(slide, Inches(8.55), Inches(1.7), Inches(3.65), Inches(3.05), [
        "Company score: weighted 0–2 score across five core metrics.",
        "Industry score: average 0–2 score across the selected industry track.",
        "Final score: 60% company score and 40% industry score.",
        "Loan brackets tighten approval thresholds as requested amount increases.",
    ], 13)
    add_footer(slide, 9)

    # 10
    slide = prs.slides.add_slide(blank)
    slides.append(slide)
    add_slide_title(slide, "Industry-Specific Underwriting", "Risk segmentation")
    industries = [
        ("Tech", "growth, margin, concentration, burn coverage, NRR", BLUE),
        ("Hospitality", "RevPAR, GOP/room, occupancy, cap rate, current ratio", AMBER),
        ("Retail", "sales/sqft, GMROI, inventory turns, gross margin, same-store sales", CYAN),
        ("Healthcare", "operating margin, cash on hand, AR days, payer mix, collections", EMERALD),
        ("Industrials", "asset turnover, interest coverage, margin, OCF, backlog", VIOLET),
    ]
    for i, (name, detail, color) in enumerate(industries):
        x = Inches(0.74 + (i % 3) * 4.05)
        y = Inches(1.78 + (i // 3) * 1.72)
        add_card(slide, x, y, Inches(3.55), Inches(1.22), name, detail, color)
    add_textbox(slide, Inches(0.78), Inches(5.52), Inches(11.2), Inches(0.5), "The important design choice is traceability: each industry metric produces a tier label, score, and triggered-rule payload that can be displayed to analysts.", 15, True, NAVY)
    add_footer(slide, 10)

    # 11
    slide = prs.slides.add_slide(blank)
    slides.append(slide)
    add_slide_title(slide, "ML Default Probability", "Risk model")
    add_textbox(slide, Inches(0.75), Inches(1.72), Inches(4.8), Inches(0.62), "The ML layer is designed as an additional signal, not a black-box replacement for policy.", 19, True, NAVY)
    add_bullets(slide, Inches(0.78), Inches(2.56), Inches(4.72), Inches(2.2), [
        "Training script supports LendingClub CSV input or synthetic fallback data.",
        "Feature set mirrors the application schema plus engineered ratios.",
        "Model pipeline combines logistic regression and gradient boosting through soft voting.",
        "API stores probability on the existing decision and logs the scoring event.",
    ], 13)
    add_image_placeholder(slide, Inches(6.05), Inches(1.78), Inches(5.85), Inches(4.4), "{{ INSERT_IMAGE: ROC curve / feature importance / model card }}")
    add_footer(slide, 11)

    # 12
    slide = prs.slides.add_slide(blank)
    slides.append(slide)
    add_slide_title(slide, "LLM-Assisted Operations", "Automation")
    add_card(slide, Inches(0.74), Inches(1.78), Inches(3.65), Inches(1.22), "Document Extraction", "Financial statements and tax returns route to `/llm/extract/{doc_id}` for structured field extraction.", BLUE)
    add_card(slide, Inches(0.74), Inches(3.22), Inches(3.65), Inches(1.22), "Decision Explanation", "Decision outcomes can be translated into plain-English rationale using the same triggered-rule payload.", EMERALD)
    add_card(slide, Inches(0.74), Inches(4.66), Inches(3.65), Inches(1.22), "Provider Abstraction", "OpenAI and Anthropic are swappable through `LLM_PROVIDER` and provider-specific API keys.", CYAN)
    add_image_placeholder(slide, Inches(5.05), Inches(1.78), Inches(6.88), Inches(4.12), "{{ INSERT_SCREENSHOT: Analyst explanation panel }}")
    add_footer(slide, 12)

    # 13
    slide = prs.slides.add_slide(blank)
    slides.append(slide)
    add_slide_title(slide, "Underwriter Dashboard", "Portfolio view")
    add_bullets(slide, Inches(0.78), Inches(1.78), Inches(4.75), Inches(2.46), [
        "KPI cards summarize exposure, active deals, approval rate, and average loan size.",
        "Charts show monthly originations and status breakdown.",
        "Recent applications link directly into the credit file.",
        "Application table supports search, status filter, sorting, and pagination.",
    ], 14)
    add_image_placeholder(slide, Inches(6.03), Inches(1.68), Inches(5.98), Inches(4.82), "{{ INSERT_SCREENSHOT: Dashboard overview and applications table }}")
    add_footer(slide, 13)

    # 14
    slide = prs.slides.add_slide(blank)
    slides.append(slide)
    add_slide_title(slide, "Decision Detail & Audit Trail", "Analyst workspace")
    shot = screenshot_path()
    if shot:
        add_picture_fit(slide, shot, Inches(0.78), Inches(1.62), Inches(4.65), Inches(4.85))
    else:
        add_image_placeholder(slide, Inches(0.78), Inches(1.62), Inches(4.65), Inches(4.85), "{{ INSERT_SCREENSHOT: Decision detail panel }}")
    add_bullets(slide, Inches(5.85), Inches(1.78), Inches(5.82), Inches(3.4), [
        "Application detail page combines loan request, borrower profile, financial ratios, radar chart, decision panel, and audit log.",
        "Analysts can run evaluation, approve, decline, flag for review, or mark approved deals as funded.",
        "Every material event is written to `audit_log` for traceability.",
    ], 15)
    add_footer(slide, 14)

    # 15
    slide = prs.slides.add_slide(blank)
    slides.append(slide)
    add_slide_title(slide, "Roadmap & Demo Narrative", "Close")
    add_textbox(slide, Inches(0.75), Inches(1.72), Inches(4.9), Inches(0.65), "The current build is a credible end-to-end platform foundation. The next work should harden production readiness and decision quality.", 19, True, NAVY)
    add_bullets(slide, Inches(0.78), Inches(2.58), Inches(4.7), Inches(2.4), [
        "Seed realistic portfolio data for dashboard demos.",
        "Persist uploaded files in local/dev storage before Supabase handoff.",
        "Add analyst-facing model card and validation metrics.",
        "Expand test coverage around policy thresholds and edge cases.",
    ], 14)
    add_card(slide, Inches(6.1), Inches(1.76), Inches(5.45), Inches(1.04), "Recommended Demo Path", "Landing page → application steps → decision evaluation → application detail → pipeline board → audit log.", BLUE)
    add_card(slide, Inches(6.1), Inches(3.08), Inches(5.45), Inches(1.04), "Investor / Judge Message", "The platform demonstrates workflow depth, policy explainability, and a path from prototype to regulated lending operations.", EMERALD)
    add_image_placeholder(slide, Inches(6.1), Inches(4.52), Inches(5.45), Inches(1.25), "{{ INSERT_FINAL_SCREENSHOT_OR_TEAM_LOGO }}")
    add_footer(slide, 15)

    OUT_DIR.mkdir(exist_ok=True)
    prs.save(DECK_PATH)
    print(DECK_PATH)


if __name__ == "__main__":
    build_deck()

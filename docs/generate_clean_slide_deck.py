from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "LoanApproval_Product_Deck_clean.pptx"
ASSETS = ROOT / "docs" / "deck_assets"
SHOTS = ROOT / "docs" / "deck_screenshots"

NAVY = RGBColor(15, 23, 42)
SLATE = RGBColor(51, 65, 85)
MUTED = RGBColor(100, 116, 139)
LINE = RGBColor(226, 232, 240)
SOFT = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
BLUE = RGBColor(29, 78, 216)
CYAN = RGBColor(8, 145, 178)
GREEN = RGBColor(5, 150, 105)
AMBER = RGBColor(217, 119, 6)
RED = RGBColor(220, 38, 38)
VIOLET = RGBColor(109, 40, 217)

W = Inches(13.333)
H = Inches(7.5)


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:], 16))


def set_text(run, size: int, color=SLATE, bold=False):
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def text(slide, x, y, w, h, value: str, size=14, color=SLATE, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for side in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, side, 0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = value
    set_text(r, size, color, bold)
    return shape


def bullet_list(slide, x, y, w, h, items: list[str], size=13):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(7)
        r = p.add_run()
        r.text = f"- {item}"
        set_text(r, size, SLATE, False)
    return box


def line(slide, y):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), y, Inches(12.08), Inches(0.01))
    s.fill.solid()
    s.fill.fore_color.rgb = LINE
    s.line.fill.background()


def title(slide, kicker: str, heading: str, n: int):
    text(slide, Inches(0.62), Inches(0.38), Inches(3.5), Inches(0.22), kicker.upper(), 8, BLUE, True)
    text(slide, Inches(0.62), Inches(0.68), Inches(8.8), Inches(0.48), heading, 24, NAVY, True)
    line(slide, Inches(1.28))
    text(slide, Inches(0.62), Inches(7.03), Inches(2.2), Inches(0.18), "LoanApproval", 8, MUTED, True)
    text(slide, Inches(12.0), Inches(7.03), Inches(0.7), Inches(0.18), f"{n:02d}", 8, MUTED, False, PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, heading: str, body: str, accent=BLUE):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = LINE
    s.line.width = Pt(1)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Inches(0.06), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    text(slide, x + Inches(0.22), y + Inches(0.18), w - Inches(0.44), Inches(0.26), heading, 12, NAVY, True)
    text(slide, x + Inches(0.22), y + Inches(0.56), w - Inches(0.44), h - Inches(0.66), body, 10, SLATE)
    return s


def placeholder(slide, x, y, w, h, label: str):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = SOFT
    s.line.color.rgb = LINE
    s.line.width = Pt(1)
    tf = s.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    set_text(r, 10, MUTED, True)
    return s


def stat(slide, x, y, label: str, value: str, color=BLUE):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, Inches(2.35), Inches(0.78))
    s.fill.solid()
    s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = rgb("#334155")
    s.line.width = Pt(0.75)
    text(slide, x + Inches(0.16), y + Inches(0.12), Inches(2.0), Inches(0.16), label.upper(), 7, rgb("#cbd5e1"), True)
    text(slide, x + Inches(0.16), y + Inches(0.36), Inches(2.0), Inches(0.28), value, 14, color, True)


def image_fit(slide, path: Path, x, y, w, h):
    if not path.exists():
        placeholder(slide, x, y, w, h, f"Missing image: {path.name}")
        return
    with Image.open(path) as img:
        iw, ih = img.size
    scale = min(w / iw, h / ih)
    pic_w = int(iw * scale)
    pic_h = int(ih * scale)
    slide.shapes.add_picture(str(path), x + (w - pic_w) / 2, y + (h - pic_h) / 2, width=pic_w, height=pic_h)


def shot_path(name: str) -> Path:
    return SHOTS / name


def screenshot() -> Path | None:
    matches = sorted(ROOT.glob("Screenshot 2026-04-22*PM.png"))
    return matches[0] if matches else None


def add_flow(slide, x, y, labels: list[tuple[str, str, RGBColor]]):
    gap = Inches(0.18)
    box_w = Inches(2.15)
    for i, (head, body, color) in enumerate(labels):
        bx = x + i * (box_w + gap)
        card(slide, bx, y, box_w, Inches(1.18), head, body, color)
        if i < len(labels) - 1:
            text(slide, bx + box_w + Inches(0.04), y + Inches(0.44), Inches(0.12), Inches(0.2), ">", 12, MUTED, True)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    # 1
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    text(s, Inches(0.78), Inches(0.68), Inches(4.2), Inches(0.24), "CORNELL FINTECH CLUB | SPRING 2026", 9, rgb("#93c5fd"), True)
    text(s, Inches(0.78), Inches(1.48), Inches(6.8), Inches(0.72), "LoanApproval", 40, WHITE, True)
    text(s, Inches(0.80), Inches(2.35), Inches(6.55), Inches(0.7), "A full-stack lending platform for borrower intake, credit policy automation, document workflows, and underwriter review.", 17, rgb("#cbd5e1"))
    stat(s, Inches(0.82), Inches(4.12), "Lifecycle", "Apply to fund", CYAN)
    stat(s, Inches(3.42), Inches(4.12), "Decisioning", "Rules + ML + LLM", GREEN)
    stat(s, Inches(6.02), Inches(4.12), "Stack", "Next.js + FastAPI", rgb("#93c5fd"))
    image_fit(s, shot_path("dashboard.png"), Inches(9.03), Inches(1.22), Inches(3.45), Inches(4.55))
    text(s, Inches(0.78), Inches(7.03), Inches(2.2), Inches(0.18), "LoanApproval", 8, rgb("#94a3b8"), True)
    text(s, Inches(12.0), Inches(7.03), Inches(0.7), Inches(0.18), "01", 8, rgb("#94a3b8"), False, PP_ALIGN.RIGHT)

    # 2
    s = prs.slides.add_slide(blank)
    title(s, "Core thesis", "Product At A Glance", 2)
    text(s, Inches(0.72), Inches(1.72), Inches(5.2), Inches(0.7), "LoanApproval converts fragmented SMB lending workflows into a structured credit file and a transparent underwriting decision.", 20, NAVY, True)
    bullet_list(s, Inches(0.75), Inches(2.76), Inches(5.1), Inches(2.15), [
        "Borrower-facing application flow collects business, loan, financial, and document data.",
        "Underwriter dashboard turns submissions into a searchable, status-driven pipeline.",
        "Credit policy is deterministic first, with ML and LLM layers supporting risk review and productivity.",
    ])
    card(s, Inches(6.35), Inches(1.72), Inches(2.75), Inches(1.08), "Borrower", "Guided application, document staging, confirmation, and status lookup.", BLUE)
    card(s, Inches(9.42), Inches(1.72), Inches(2.75), Inches(1.08), "Underwriter", "Portfolio overview, pipeline, queue, detail page, and action buttons.", GREEN)
    card(s, Inches(6.35), Inches(3.14), Inches(2.75), Inches(1.08), "Policy", "Weighted ratios, industry tracks, loan brackets, hard stops, and caps.", AMBER)
    card(s, Inches(9.42), Inches(3.14), Inches(2.75), Inches(1.08), "Platform", "FastAPI, PostgreSQL, Supabase Storage, OpenAI or Anthropic.", CYAN)
    placeholder(s, Inches(6.35), Inches(4.72), Inches(5.82), Inches(1.18), "INSERT OVERVIEW COMPOSITE")

    # 3
    s = prs.slides.add_slide(blank)
    title(s, "Workflow", "End-to-End Loan Lifecycle", 3)
    add_flow(s, Inches(0.72), Inches(1.8), [
        ("Apply", "Borrower profile and loan request", BLUE),
        ("Financials", "Core values and calculated ratios", CYAN),
        ("Documents", "Statements, returns, bank records", GREEN),
        ("Evaluate", "Rules engine and model scoring", AMBER),
        ("Manage", "Pipeline, actions, audit log", VIOLET),
    ])
    text(s, Inches(0.78), Inches(3.6), Inches(10.6), Inches(0.35), "Application state model", 16, NAVY, True)
    add_flow(s, Inches(0.78), Inches(4.18), [
        ("New", "Submitted and awaiting review", MUTED),
        ("In Review", "Manual review or capped decision", AMBER),
        ("Approved", "Credit decision accepted", GREEN),
        ("Funded", "Post-approval funding state", VIOLET),
        ("Declined", "Terminal rejection branch", RED),
    ])

    # 4
    s = prs.slides.add_slide(blank)
    title(s, "Frontend", "Borrower Entry Point", 4)
    text(s, Inches(0.76), Inches(1.72), Inches(4.75), Inches(0.52), "The landing page is simple: explain the offer, set expectations, and route directly into the application.", 18, NAVY, True)
    bullet_list(s, Inches(0.78), Inches(2.55), Inches(4.65), Inches(2.05), [
        "CTA points to the first application step.",
        "Loan products cover term loan, line of credit, equipment finance, and bridge loan.",
        "Eligibility copy reduces low-quality submissions before intake begins.",
    ])
    image_fit(s, shot_path("landing.png"), Inches(6.0), Inches(1.65), Inches(6.1), Inches(4.8))

    # 5
    s = prs.slides.add_slide(blank)
    title(s, "Borrower flow", "Guided Application Intake", 5)
    card(s, Inches(0.75), Inches(1.78), Inches(3.55), Inches(1.05), "Business Profile", "Company name, DBA, industry, legal structure, years in operation, EIN, and address.", BLUE)
    card(s, Inches(0.75), Inches(3.1), Inches(3.55), Inches(1.05), "Loan Request", "Amount, purpose, optional details, and desired term.", CYAN)
    card(s, Inches(0.75), Inches(4.42), Inches(3.55), Inches(1.05), "Review and Submit", "Certified payload creates borrower and application records.", GREEN)
    image_fit(s, shot_path("apply_step1.png"), Inches(5.05), Inches(1.65), Inches(6.95), Inches(4.82))

    # 6
    s = prs.slides.add_slide(blank)
    title(s, "Credit file", "Financial Data Capture", 6)
    text(s, Inches(0.76), Inches(1.72), Inches(4.9), Inches(0.58), "Structured financial data gives the backend enough information to calculate policy metrics without manual spreadsheet work.", 18, NAVY, True)
    bullet_list(s, Inches(0.78), Inches(2.58), Inches(4.8), Inches(2.25), [
        "Required fields include revenue, EBITDA, debt, assets, collateral, debt service, and bankruptcies.",
        "Optional fields support current ratio and interest coverage.",
        "The frontend calculates DSCR, LTV, and leverage while the borrower types.",
    ])
    card(s, Inches(6.05), Inches(1.78), Inches(2.65), Inches(0.95), "DSCR", "EBITDA / annual debt service", GREEN)
    card(s, Inches(9.0), Inches(1.78), Inches(2.65), Inches(0.95), "LTV", "loan amount / collateral", AMBER)
    card(s, Inches(6.05), Inches(3.06), Inches(2.65), Inches(0.95), "Leverage", "existing debt / EBITDA", BLUE)
    card(s, Inches(9.0), Inches(3.06), Inches(2.65), Inches(0.95), "Liquidity", "current assets / current liabilities", CYAN)
    image_fit(s, shot_path("apply_step3.png"), Inches(6.05), Inches(4.38), Inches(5.6), Inches(1.55))

    # 7
    s = prs.slides.add_slide(blank)
    title(s, "Evidence", "Document Workflow", 7)
    text(s, Inches(0.76), Inches(1.72), Inches(4.8), Inches(0.52), "Document handling is modeled as part of the application record, not as an external email process.", 18, NAVY, True)
    bullet_list(s, Inches(0.78), Inches(2.55), Inches(4.75), Inches(2.05), [
        "Required: financial statements, tax returns, and bank statements.",
        "Optional: accounts receivable aging and collateral documents.",
        "Upload records metadata, storage path, and an audit event.",
    ])
    image_fit(s, shot_path("documents.png"), Inches(6.0), Inches(1.65), Inches(6.1), Inches(4.82))

    # 8
    s = prs.slides.add_slide(blank)
    title(s, "Platform", "Backend Architecture", 8)
    image_fit(s, ASSETS / "architecture.png", Inches(0.75), Inches(1.55), Inches(11.84), Inches(4.95))

    # 9
    s = prs.slides.add_slide(blank)
    title(s, "Decisioning", "Credit Policy Engine", 9)
    image_fit(s, ASSETS / "decision_model.png", Inches(0.72), Inches(1.48), Inches(7.25), Inches(4.88))
    bullet_list(s, Inches(8.35), Inches(1.82), Inches(3.8), Inches(3.0), [
        "Company score: five weighted core metrics on a 0 to 2 scale.",
        "Industry score: five metrics for the selected industry track.",
        "Final score: 60 percent company, 40 percent industry.",
        "Loan brackets raise approval thresholds for larger requests.",
    ])

    # 10
    s = prs.slides.add_slide(blank)
    title(s, "BA placeholder", "Manual Credit Formula Workbench", 10)
    text(s, Inches(0.76), Inches(1.72), Inches(6.2), Inches(0.42), "Blank working slide for business-analysis notes behind the manual credit formulas and analyst-facing rationale.", 17, NAVY, True)
    text(s, Inches(0.78), Inches(2.36), Inches(11.0), Inches(0.28), "Rough outline to fill in later", 13, MUTED, True)
    card(s, Inches(0.78), Inches(2.86), Inches(2.72), Inches(1.18), "Formula name", "Example: DSCR, LTV, leverage, industry metric, hard stop, or cap.", BLUE)
    card(s, Inches(3.78), Inches(2.86), Inches(2.72), Inches(1.18), "Manual definition", "Plain-English BA definition and any source notes.", CYAN)
    card(s, Inches(6.78), Inches(2.86), Inches(2.72), Inches(1.18), "Inputs", "Required fields, fallback logic, and missing-data handling.", GREEN)
    card(s, Inches(9.78), Inches(2.86), Inches(2.72), Inches(1.18), "Thresholds", "Approve / review / concern / reject tiers by bracket or segment.", AMBER)
    card(s, Inches(0.78), Inches(4.42), Inches(3.62), Inches(1.25), "Decision impact", "Does this affect company score, industry score, final score, hard stop, or review cap?", VIOLET)
    card(s, Inches(4.72), Inches(4.42), Inches(3.62), Inches(1.25), "Applicant explanation", "How the rule should be described in the LLM-generated explanation.", BLUE)
    card(s, Inches(8.66), Inches(4.42), Inches(3.62), Inches(1.25), "Validation notes", "Test cases, edge cases, owner, and approval status.", RED)

    # 11
    s = prs.slides.add_slide(blank)
    title(s, "Risk segmentation", "Industry-Specific Underwriting", 11)
    rows = [
        ("Tech", "growth, gross margin, concentration, burn coverage, NRR", BLUE),
        ("Hospitality", "RevPAR, GOP per room, occupancy, cap rate, current ratio", AMBER),
        ("Retail", "sales per square foot, GMROI, inventory turns, gross margin", CYAN),
        ("Healthcare", "operating margin, cash on hand, AR days, payer mix", GREEN),
        ("Industrials", "asset turnover, interest coverage, margin, OCF, backlog", VIOLET),
    ]
    for i, row in enumerate(rows):
        card(s, Inches(0.78 + (i % 3) * 4.05), Inches(1.75 + (i // 3) * 1.55), Inches(3.55), Inches(1.06), *row)
    text(s, Inches(0.82), Inches(5.4), Inches(10.9), Inches(0.48), "Each metric produces a tier label, numeric score, and triggered-rule payload so analysts can see exactly why the policy engine reached its result.", 15, NAVY, True)

    # 12
    s = prs.slides.add_slide(blank)
    title(s, "Risk model", "ML Default Probability", 12)
    text(s, Inches(0.76), Inches(1.72), Inches(4.95), Inches(0.52), "The model score is a supplemental risk signal. It does not replace policy thresholds or hard stops.", 18, NAVY, True)
    bullet_list(s, Inches(0.78), Inches(2.52), Inches(4.8), Inches(2.45), [
        "Training supports LendingClub input or synthetic fallback data.",
        "Features mirror application fields plus debt-to-equity, DSCR, and revenue-to-loan.",
        "Pipeline combines logistic regression and gradient boosting through soft voting.",
        "API writes probability to the decision record and logs the scoring event.",
    ])
    placeholder(s, Inches(6.05), Inches(1.72), Inches(5.95), Inches(4.55), "INSERT IMAGE: MODEL CARD OR ROC CURVE")

    # 13
    s = prs.slides.add_slide(blank)
    title(s, "Automation", "LLM-Assisted Operations", 13)
    card(s, Inches(0.78), Inches(1.78), Inches(3.55), Inches(1.1), "Extraction", "Financial statements and tax returns route to the LLM extraction endpoint.", BLUE)
    card(s, Inches(0.78), Inches(3.16), Inches(3.55), Inches(1.1), "Explanation", "Triggered rules and score context become plain-English rationale.", GREEN)
    card(s, Inches(0.78), Inches(4.54), Inches(3.55), Inches(1.1), "Provider Choice", "OpenAI and Anthropic are swappable through environment configuration.", CYAN)
    placeholder(s, Inches(5.08), Inches(1.78), Inches(6.95), Inches(4.28), "INSERT SCREENSHOT: ANALYST EXPLANATION")

    # 14
    s = prs.slides.add_slide(blank)
    title(s, "Portfolio view", "Underwriter Dashboard", 14)
    bullet_list(s, Inches(0.78), Inches(1.82), Inches(4.8), Inches(2.72), [
        "KPI cards summarize exposure, active deals, approval rate, and average loan size.",
        "Charts show monthly originations and status breakdown.",
        "Recent applications link directly to the credit file.",
        "Application table supports search, status filtering, sorting, and pagination.",
    ])
    image_fit(s, shot_path("dashboard.png"), Inches(6.02), Inches(1.62), Inches(6.05), Inches(4.92))

    # 15
    s = prs.slides.add_slide(blank)
    title(s, "Analyst workspace", "Decision Detail and Audit Trail", 15)
    shot = screenshot()
    if shot:
        s.shapes.add_picture(str(shot), Inches(0.8), Inches(1.58), height=Inches(4.92))
    else:
        placeholder(s, Inches(0.8), Inches(1.58), Inches(4.62), Inches(4.92), "INSERT SCREENSHOT: DECISION PANEL")
    bullet_list(s, Inches(5.82), Inches(1.78), Inches(5.85), Inches(2.7), [
        "Detail page combines loan request, borrower profile, financial ratios, decision panel, and audit log.",
        "Analysts can evaluate, approve, decline, flag for review, or mark approved deals as funded.",
        "Audit events preserve the application lifecycle and decision rationale.",
    ], 14)

    # 16
    s = prs.slides.add_slide(blank)
    title(s, "Close", "Roadmap and Demo Narrative", 16)
    text(s, Inches(0.76), Inches(1.72), Inches(5.05), Inches(0.58), "The current build demonstrates a credible end-to-end lending workflow. Next steps should focus on production hardening and richer evidence.", 18, NAVY, True)
    bullet_list(s, Inches(0.78), Inches(2.58), Inches(4.85), Inches(2.4), [
        "Seed realistic portfolio data for demos.",
        "Harden local and production document storage.",
        "Add model validation metrics and an analyst-facing model card.",
        "Expand tests around threshold boundaries and policy edge cases.",
    ])
    card(s, Inches(6.12), Inches(1.78), Inches(5.55), Inches(1.05), "Demo Path", "Landing page -> application steps -> evaluation -> decision detail -> pipeline -> audit log.", BLUE)
    card(s, Inches(6.12), Inches(3.18), Inches(5.55), Inches(1.05), "Takeaway", "Workflow depth, policy explainability, and an extensible path toward regulated lending operations.", GREEN)
    placeholder(s, Inches(6.12), Inches(4.68), Inches(5.55), Inches(1.15), "INSERT FINAL SCREENSHOT OR TEAM LOGO")

    OUT.parent.mkdir(exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
